from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
TEAL = RGBColor(0x02, 0x80, 0x90)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x20, 0x25, 0x2E)
LIGHT_BG = RGBColor(0xF6, 0xF8, 0xFC)
GREEN = RGBColor(0x2C, 0xA5, 0x8D)
ORANGE = RGBColor(0xF9, 0x61, 0x67)


def add_bg(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title(slide, text: str, subtitle: str | None = None, dark: bool = False) -> None:
    color = WHITE if dark else NAVY
    sub_color = ICE if dark else DARK
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = sub_color
        p2.font.name = "Calibri"


def card(slide, x: float, y: float, w: float, h: float, title: str, body: str) -> None:
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = ICE
    tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.2), Inches(w - 0.5), Inches(h - 0.35))
    tf = tb.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK


def build_presentation(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s1, NAVY)
    add_title(
        s1,
        "SkylineA2A Project Overview",
        "Foundry + Microsoft Agent Framework + APIM monetization",
        dark=True,
    )
    pill = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.4), Inches(0.7))
    pill.fill.solid()
    pill.fill.fore_color.rgb = TEAL
    pill.line.fill.background()
    t = pill.text_frame.paragraphs[0]
    t.text = "Two A2A agents • One monetized gateway"
    t.font.bold = True
    t.font.color.rgb = WHITE
    t.font.size = Pt(20)
    t.alignment = PP_ALIGN.CENTER

    # Slide 2: What we built
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s2, LIGHT_BG)
    add_title(s2, "What we built", "Business-ready A2A demo with monetization controls")
    card(
        s2,
        0.8,
        1.8,
        5.9,
        2.25,
        "Demo 1 — Foundry Agent Service",
        "Docs Expert agent on GPT-5.2, exposed as A2A and grounded with Microsoft Learn MCP.",
    )
    card(
        s2,
        6.9,
        1.8,
        5.6,
        2.25,
        "Demo 2 — Microsoft Agent Framework",
        "DevOps Helper agent on GPT-5.2, exposed as A2A from a Python server (local or ACA).",
    )
    card(
        s2,
        0.8,
        4.3,
        11.7,
        2.1,
        "Unified front door — Azure API Management",
        "Products, subscriptions, quotas and policy-based governance applied consistently to both agent runtimes.",
    )

    # Slide 3: Architecture view
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s3, WHITE)
    add_title(s3, "Solution architecture", "A2A clients consume agents through APIM")

    client = s3.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.5), Inches(2.2), Inches(1.0))
    client.fill.solid()
    client.fill.fore_color.rgb = ICE
    client.line.color.rgb = NAVY
    client.text_frame.paragraphs[0].text = "A2A Client"
    client.text_frame.paragraphs[0].font.bold = True
    client.text_frame.paragraphs[0].font.color.rgb = NAVY
    client.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    apim = s3.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(1.8), Inches(6.2), Inches(2.4))
    apim.fill.solid()
    apim.fill.fore_color.rgb = NAVY
    apim.line.fill.background()
    apim_tf = apim.text_frame
    apim_tf.clear()
    p = apim_tf.paragraphs[0]
    p.text = "Azure API Management"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = apim_tf.add_paragraph()
    p2.text = "Products: Free / Pro\nPolicies: key auth, quota/rate, metering"
    p2.font.size = Pt(14)
    p2.font.color.rgb = ICE
    p2.alignment = PP_ALIGN.CENTER

    d1 = s3.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.8), Inches(4.0), Inches(1.6))
    d1.fill.solid()
    d1.fill.fore_color.rgb = RGBColor(0xEC, 0xF7, 0xF9)
    d1.line.color.rgb = TEAL
    d1.text_frame.paragraphs[0].text = "Demo 1\nFoundry Agent Service\nA2A endpoint"
    d1.text_frame.paragraphs[0].font.bold = True
    d1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    d2 = s3.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.7), Inches(4.8), Inches(4.0), Inches(1.6))
    d2.fill.solid()
    d2.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xFF)
    d2.line.color.rgb = NAVY
    d2.text_frame.paragraphs[0].text = "Demo 2\nMAF A2A server\non ACA"
    d2.text_frame.paragraphs[0].font.bold = True
    d2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 4: APIM monetization strategy
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s4, LIGHT_BG)
    add_title(s4, "API Management monetization strategy", "Commercial control without changing agent code")
    card(
        s4,
        0.8,
        1.8,
        4.0,
        2.2,
        "Products",
        "Free and Pro products define traffic envelopes and package tiers for customers.",
    )
    card(
        s4,
        4.95,
        1.8,
        4.0,
        2.2,
        "Subscriptions",
        "Each consumer gets a subscription key that maps usage to a plan and tenant.",
    )
    card(
        s4,
        9.1,
        1.8,
        3.4,
        2.2,
        "Policies",
        "rate-limit, quota and telemetry applied centrally for both A2A APIs.",
    )
    good = s4.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.45), Inches(5.85), Inches(1.75))
    good.fill.solid()
    good.fill.fore_color.rgb = RGBColor(0xE7, 0xF6, 0xEF)
    good.line.color.rgb = GREEN
    gt = good.text_frame
    gt.text = "✅ Implemented now\n• Product-tier access\n• Key-based identity\n• Quota/rate controls\n• Per-call metering signals"
    gt.paragraphs[0].font.bold = True
    gt.paragraphs[0].font.color.rgb = GREEN
    for p in gt.paragraphs:
        p.font.size = Pt(14)

    next_box = s4.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(4.45), Inches(5.65), Inches(1.75))
    next_box.fill.solid()
    next_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xEE)
    next_box.line.color.rgb = ORANGE
    nt = next_box.text_frame
    nt.text = "⚠️ Next step\nToken-level billing can be layered with custom telemetry + external metering."
    nt.paragraphs[0].font.bold = True
    nt.paragraphs[0].font.color.rgb = ORANGE
    for p in nt.paragraphs:
        p.font.size = Pt(14)

    # Slide 5: Final message
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s5, NAVY)
    add_title(
        s5,
        "Result: a credible long-term A2A monetization blueprint",
        "Reusable for future agents with minimal integration overhead",
        dark=True,
    )
    summary = s5.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.8), Inches(3.0))
    tf = summary.text_frame
    tf.clear()
    for line in [
        "• Same APIM monetization model for Foundry-hosted and MAF-hosted agents",
        "• Clear product strategy (Free vs Pro) and measurable consumption boundaries",
        "• End-to-end validation already working through APIM routes",
        "• Ready to evolve toward enterprise billing and governance patterns",
    ]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.level = 0

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    build_presentation(Path(__file__).resolve().parent / "SkylineA2A_Project_Overview.pptx")
